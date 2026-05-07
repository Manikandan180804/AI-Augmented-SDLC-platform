from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def set_slide_background(slide, image_path):
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, width=Presentation().slide_width, height=Presentation().slide_height)

def add_title_with_styling(slide, text, font_size=44, color=RGBColor(255, 255, 255)):
    title = slide.shapes.title
    title.text = text
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = color
        paragraph.alignment = PP_ALIGN.LEFT

def create_presentation():
    prs = Presentation()
    
    # Standard Slide Dimensions (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Define Layouts
    # 0: Title, 1: Title & Content, 6: Blank
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]

    bg_path = 'background.png'

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    set_slide_background(slide, bg_path)
    
    title = slide.shapes.title
    title.text = "AI-Augmented Software Development Lifecycle (SDLC) Engine"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Automating Efficiency from Requirements to Production\n\nPresenter: [Your Name]\nCompany: Prodapt (Internship 2026)"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    # --- Slide 2: The Problem Statement ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    # Background color (Dark)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 25, 47) # Deep Navy

    title = slide.shapes.title
    title.text = "The Problem Statement (The 'Why')"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 183, 235) # Cyan

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Current SDLC Challenges:"
    
    points = [
        ("Ambiguity", "Vague requirements lead to 40% of project rework."),
        ("Bottlenecks", "Manual code reviews slow down the 'Lead Time to Production'."),
        ("Gaps", "Inconsistent test coverage and missing traceability."),
        ("Risk", "Blind deployments lead to production incidents and developer burnout.")
    ]
    
    for bold_text, desc in points:
        p = tf.add_paragraph()
        p.text = f"{bold_text}: {desc}"
        p.level = 0
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(20)

    # Add TL Quote
    p = tf.add_paragraph()
    p.text = "" # Spacer
    p = tf.add_paragraph()
    p.text = 'TL Quote: "There is a critical need for intelligent automation that catches errors before they reach production."'
    p.font.italic = True
    p.font.color.rgb = RGBColor(255, 165, 0) # Orange accent
    p.font.size = Pt(22)

    # --- Slide 3: The Solution ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 25, 47)

    title = slide.shapes.title
    title.text = "The Solution: AI-SDLC Platform"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 183, 235)

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Vision: A unified platform that 'Augments' human intelligence."
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    items = [
        "Core: 6 specialized AI modules running on a high-performance stack.",
        "Stack: FastAPI (Backend) + React (Frontend).",
        "Goal: Zero-friction development with 100% traceability."
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(24)

    # --- Slide 4: Architecture ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 25, 47)

    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 183, 235)

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    arch_items = [
        "Frontend: React 18 with Glassmorphism Design System",
        "Backend: FastAPI for asynchronous AI orchestration",
        "Knowledge: RAG Engine with FAISS Vector Store for codebase context",
        "Connectivity: Event-driven Webhooks (GitHub / Jira)",
        "AI Engine: OpenAI GPT-4o (Structured Output API)"
    ]
    for item in arch_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(22)

    # --- Slides 5-10: Modules ---
    modules = [
        ("Module 1: Requirements Analyzer", [
            "Detects ambiguities in natural language",
            "Automatically drafts Gherkin User Stories",
            "Flags Scope Creep and Undefined Dependencies",
            "Calculates a Completeness Score (0-100)"
        ], "Value: Replaces vague talk with testable requirements."),
        ("Module 2: AI Code Reviewer", [
            "Security-first architectural review",
            "Identifies OWASP Top 10 vulnerabilities",
            "Provides 'Before vs. After' code diffs with suggested fixes",
            "Checks for missing unit tests for the changed code"
        ], "Value: Shields production from security flaws."),
        ("Module 3: Smart Test Generator", [
            "Generates full pytest / jest files from raw function bodies",
            "Covers edge cases (null values, boundary limits)",
            "Prioritizes tests based on 'Regression Risk'"
        ], "Value: Ensures high-quality coverage without manual overhead."),
        ("Module 4: Deployment Risk Scorer", [
            "Calculates a 0-100 risk score based on change signals",
            "Escalates Auth changes and DB migrations",
            "Generates Rollback Plan and Monitoring Checklist"
        ], "Value: Minimizes downtime and deployment anxiety."),
        ("Module 5: Traceability Engine", [
            "Maps: Requirement → Story → Commit → Test → Deployment",
            "Detects 'Orphaned Code' (Code with no requirement)",
            "Visually traces the 'Chain of Custody'"
        ], "Value: Simplifies audits and ensures business alignment."),
        ("Module 6: Multi-Agent Bug Fixer", [
            "Collaborative workflow: Analyst → Researcher → Fixer → QA",
            "Automated root cause analysis from bug reports",
            "Retrieves codebase context using RAG",
            "Generates verified code fixes with regression warnings"
        ], "Value: Slashes 'Mean Time To Repair' (MTTR) by 70%.")
    ]

    for m_title, m_points, m_value in modules:
        slide = prs.slides.add_slide(bullet_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(10, 25, 47)
        
        title = slide.shapes.title
        title.text = m_title
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 183, 235)
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        for pt in m_points:
            p = tf.add_paragraph()
            p.text = pt
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(20)
            p.level = 0
            
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = m_value
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 255, 127) # Spring Green
        p.font.size = Pt(22)

    # --- Slide 11: RAG & Contextual Intelligence ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 25, 47)

    title = slide.shapes.title
    title.text = "RAG & Contextual Intelligence"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 183, 235)

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    rag_points = [
        "Problem: AI often lacks context of the existing codebase.",
        "Solution: Retrieval-Augmented Generation (RAG).",
        "Implementation: Codebase is indexed into a FAISS Vector Store.",
        "Benefit: AI makes decisions based on actual project patterns, not just generic knowledge."
    ]
    for item in rag_points:
        p = tf.add_paragraph()
        p.text = item
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(22)

    # --- Slide 12: Enterprise Features ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 25, 47)

    title = slide.shapes.title
    title.text = "Enterprise Features (The 'Extra Mile')"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 183, 235)

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    features = [
        "Real-time Integration: Live GitHub/Jira Webhook feed",
        "Security: API Key Guarding and Header-based authentication",
        "Exportability: One-click JSON/PDF/PPTX reports",
        "Scalability: Asynchronous task processing with FastAPI"
    ]
    for item in features:
        p = tf.add_paragraph()
        p.text = item
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(24)

    # --- Slide 13: Roadmap ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 25, 47)

    title = slide.shapes.title
    title.text = "Future Roadmap"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 183, 235)

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    roadmap = [
        "Phase 4: Autonomous Pull Request Creation & Management",
        "Phase 5: Self-Healing Production (Auto-rollback on metric drift)",
        "Phase 6: Multi-Repo Cross-Project Intelligence"
    ]
    for item in roadmap:
        p = tf.add_paragraph()
        p.text = item
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(24)
    # --- Slide 14: Conclusion ---
    slide = prs.slides.add_slide(title_slide_layout)
    set_slide_background(slide, bg_path)
    
    title = slide.shapes.title
    title.text = "Conclusion / Q&A"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    subtitle = slide.placeholders[1]
    subtitle.text = 'Summary: Moving from manual to AI-Augmented SDLC\n\n"Technology is best when it empowers developers to focus on creativity while AI handles the complexity."'
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    output_path = 'AI_SDLC_Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation generated successfully at {output_path}")

if __name__ == "__main__":
    create_presentation()

